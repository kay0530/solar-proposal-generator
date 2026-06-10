"""
utils.py — Shared python-pptx design system for all slide generators.

Design v2: "Institutional Trust Grid" (2026-06 redesign)
  - White canvas, 12-column grid, A4 landscape
  - Orange #E8490F reserved for numbers / CTA / brand ticks (<=5% area)
  - Structural navy #1F3551 for titles and table rules
  - Warm-gray support palette; no pure red, no teal, no gradient chrome
  - Single KPI number size (28pt); HERO 48pt at most twice per deck
  - Tables styled as "audited figures": white header + navy rule,
    warm zebra rows, right-aligned numerals, horizontal hairlines only
  - Charts: before-series = gray dashed, after-series = orange solid
    (grayscale-print double encoding)

Legacy constant names (C_TEAL, C_LIGHT_CYAN, ...) are kept as aliases
mapped into the new palette so unmigrated slides keep importing, but
render inside the v2 color world. Core components (add_header_bar,
add_footer, add_kpi_card, add_section_header, add_table) are
destructively rewritten — there is no way to draw the old chrome.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Canvas
# ---------------------------------------------------------------------------

SLIDE_W = Inches(11.69)   # A4 landscape width
SLIDE_H = Inches(8.27)    # A4 landscape height

MARGIN = Inches(0.6)      # page side margin (v2: widened from 0.35)

# Vertical bands
HEADER_H = Inches(1.10)           # white header zone (eyebrow + title + rule)
CONTENT_TOP = Inches(1.30)
CONTENT_BOTTOM = Inches(7.65)
FOOTER_RULE_Y = Inches(7.78)
FOOTER_H = SLIDE_H - CONTENT_BOTTOM          # legacy name, kept for imports
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP

# ---------------------------------------------------------------------------
# 12-column grid
# ---------------------------------------------------------------------------

GRID_COLS = 12
GRID_GUTTER = Inches(0.16)
_CONTENT_W = SLIDE_W - MARGIN * 2
GRID_COL_W = (_CONTENT_W - GRID_GUTTER * (GRID_COLS - 1)) // GRID_COLS


def grid_x(col: int):
    """Left x of grid column `col` (0-based)."""
    return MARGIN + (GRID_COL_W + GRID_GUTTER) * col


def grid_w(span: int):
    """Width of `span` grid columns including internal gutters."""
    return GRID_COL_W * span + GRID_GUTTER * (span - 1)


# ---------------------------------------------------------------------------
# Color palette v2
# ---------------------------------------------------------------------------

C_ORANGE      = RGBColor(0xE8, 0x49, 0x0F)   # numbers / CTA / brand tick ONLY
C_ORANGE_DARK = RGBColor(0xC5, 0x3D, 0x0A)   # negative values, 2nd chart tone
C_TINT        = RGBColor(0xFD, 0xED, 0xE6)   # highlight cells / total rows only
C_NAVY        = RGBColor(0x1F, 0x35, 0x51)   # titles, table rules, CTA band
C_NAVY_LIGHT  = RGBColor(0x5A, 0x7B, 0xA6)   # muted navy (legacy teal remap)
C_DARK        = RGBColor(0x33, 0x33, 0x33)   # body text
C_SUB         = RGBColor(0x66, 0x66, 0x66)   # labels / captions / axis
C_FAINT       = RGBColor(0x99, 0x99, 0x99)   # footer / page numbers
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL       = RGBColor(0xF7, 0xF5, 0xF2)   # warm panel (conclusion bands)
C_ZEBRA       = RGBColor(0xF2, 0xEF, 0xEA)   # table even rows
C_HAIR        = RGBColor(0xD8, 0xD4, 0xCC)   # 0.5pt hairlines / card borders
C_CHART_GRAY  = RGBColor(0xB8, 0xB2, 0xA7)   # "before" series / tertiary
C_GRID_LINE   = RGBColor(0xE5, 0xE1, 0xDA)   # chart gridlines

# Legacy aliases — keep imports alive, remap into v2 world
C_LIGHT_ORANGE = C_TINT
C_LIGHT_GRAY   = C_PANEL
C_BORDER       = C_HAIR
C_CARD_BG      = C_NAVY
C_TEAL         = C_NAVY
C_LIGHT_TEAL   = C_NAVY_LIGHT
C_LIGHT_CYAN   = C_ZEBRA
C_RED          = C_ORANGE_DARK
C_BG_SOFT      = C_PANEL

# Chart series palette (ordered)
CHART_SERIES = [C_ORANGE, C_NAVY, C_CHART_GRAY]

# ---------------------------------------------------------------------------
# Typography scale v2 (pt)
# ---------------------------------------------------------------------------

FONT_BLACK = "Meiryo"     # display — always pair with bold=True
FONT_BODY  = "Meiryo"

SIZE_HERO    = 48     # 1 per slide max; <=2 per deck (price hero, summary)
SIZE_KPI     = 28     # the ONLY standard KPI number size
SIZE_H1      = 20     # slide title (header zone, navy)
SIZE_H2      = 14     # section heading
SIZE_LEAD    = 12.5   # standfirst / conclusion line (line_spacing 1.4)
SIZE_BODY_LG = 12.5   # legacy alias of LEAD
SIZE_H3      = 12.5   # legacy alias
SIZE_BODY    = 10.5   # body default (line_spacing 1.35)
SIZE_CAPTION = 9      # labels / eyebrows / axis / table body
SIZE_TABLE   = 9
SIZE_SMALL   = 8      # notes / footer — hard floor, never go below
SIZE_XS      = 8      # legacy alias (7pt is banned in v2)

LINE_SPACING_BODY = 1.35
LINE_SPACING_LEAD = 1.4

# ---------------------------------------------------------------------------
# Spacing scale v2 — only these four gaps
# ---------------------------------------------------------------------------

GAP_IN_CARD  = Inches(0.12)   # between elements inside a card
GAP_CARD     = Inches(0.20)   # between cards
GAP_BLOCK    = Inches(0.32)   # between blocks
GAP_SECTION  = Inches(0.48)   # between sections

# Legacy aliases
GAP_XS = Inches(0.05)
GAP_SM = GAP_IN_CARD
GAP_MD = GAP_CARD
GAP_LG = GAP_BLOCK
GAP_XL = GAP_SECTION

CARD_PAD = Inches(0.15)       # fixed card inner padding
CARD_RADIUS_PT = 4.0          # single corner radius for the whole deck
ACCENT_BAR_H = Inches(0.045)  # card top accent thickness
ACCENT_BAR_INSET = Inches(0.12)

# ---------------------------------------------------------------------------
# Template helpers
# ---------------------------------------------------------------------------

def load_template(template_path: Path) -> Presentation:
    return Presentation(str(template_path))


def add_blank_slide(prs: Presentation, layout_index: int = 6):
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


# ---------------------------------------------------------------------------
# Shape primitives
# ---------------------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill_color: RGBColor,
             border_color: Optional[RGBColor] = None, border_pt: float = 0.0):
    shape = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color: RGBColor,
                     radius_pt: float = CARD_RADIUS_PT,
                     border_color: Optional[RGBColor] = None,
                     border_pt: float = 0.0):
    shape = slide.shapes.add_shape(5, int(x), int(y), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    sp_pr = shape.element.find(qn("p:spPr"))
    if sp_pr is not None:
        prstgeom = sp_pr.find(qn("a:prstGeom"))
        if prstgeom is not None:
            av_lst = prstgeom.find(qn("a:avLst"))
            if av_lst is not None:
                for gd in av_lst.findall(qn("a:gd")):
                    if gd.get("name") == "adj":
                        min_dim = min(int(w), int(h))
                        # adj = radius / min_dim scaled to 100000, capped at pill
                        frac = min(radius_pt * 12700 / min_dim * 100000, 50000)
                        gd.set("fmla", f"val {int(frac)}")
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_oval(slide, x, y, w, h, fill_color: Optional[RGBColor] = None,
             border_color: Optional[RGBColor] = None, border_pt: float = 1.0):
    """Plain ellipse. fill_color=None → transparent fill."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(w), int(h))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def set_tracking(run, spc_pt: float = 1.2):
    """Letter-spacing on a run via a:spc (units: 1/100 pt)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(spc_pt * 100)))


def add_textbox(slide, x, y, w, h, text: str,
                font_name: str = FONT_BODY,
                font_size_pt: float = SIZE_BODY,
                font_color: RGBColor = C_DARK,
                bold: bool = False,
                align: PP_ALIGN = PP_ALIGN.LEFT,
                word_wrap: bool = True,
                line_spacing: Optional[float] = None,
                tracking_pt: Optional[float] = None,
                anchor: Optional["MSO_ANCHOR"] = None):
    """Single-run textbox. v2 adds line_spacing / tracking / vertical anchor."""
    txBox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    if anchor is not None:
        tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = font_color
    run.font.bold = bold
    if tracking_pt:
        set_tracking(run, tracking_pt)
    return txBox


def add_multiline_textbox(slide, x, y, w, h, lines: list[tuple],
                          word_wrap: bool = True,
                          line_spacing: Optional[float] = None):
    """Multi-paragraph textbox.
    lines: list of (text, font_name, size_pt, color, bold, align)."""
    txBox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    for i, (text, font_name, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return txBox


def add_number_unit(slide, x, y, w, h,
                    number: str, unit: str,
                    number_size_pt: float = None,
                    unit_size_pt: float = 11,
                    number_color: RGBColor = C_ORANGE,
                    unit_color: RGBColor = C_DARK,
                    align: PP_ALIGN = PP_ALIGN.LEFT,
                    bold_unit: bool = False):
    """Number + unit as TWO RUNS in ONE textbox, bottom-anchored so the
    pair shares a baseline (Meiryo digits/kana align well enough).
    Never stack number and unit in separate boxes."""
    if number_size_pt is None:
        number_size_pt = SIZE_KPI
    txBox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    p = tf.paragraphs[0]
    p.alignment = align
    r1 = p.add_run()
    r1.text = str(number)
    r1.font.name = FONT_BLACK
    r1.font.size = Pt(number_size_pt)
    r1.font.color.rgb = number_color
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = f" {unit}"
    r2.font.name = FONT_BODY
    r2.font.size = Pt(unit_size_pt)
    r2.font.color.rgb = unit_color
    r2.font.bold = bold_unit
    return txBox


def add_image_contain(slide, x, y, w, h, image_path: Path):
    from PIL import Image as PILImage
    img = PILImage.open(str(image_path))
    img_w, img_h = img.size
    aspect = img_w / img_h
    box_aspect = w / h
    if aspect > box_aspect:
        render_w = int(w)
        render_h = int(int(w) / aspect)
        render_x = int(x)
        render_y = int(y) + (int(h) - render_h) // 2
    else:
        render_h = int(h)
        render_w = int(int(h) * aspect)
        render_x = int(x) + (int(w) - render_w) // 2
        render_y = int(y)
    return slide.shapes.add_picture(str(image_path), render_x, render_y,
                                    render_w, render_h)


def add_line(slide, x1, y1, x2, y2, color: RGBColor, width_pt: float = 1.0,
             dash: Optional[str] = None):
    connector = slide.shapes.add_connector(1, int(x1), int(y1), int(x2), int(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    if dash:
        try:
            ln = connector.line._get_or_add_ln()
            d = ln.find(qn("a:prstDash"))
            if d is None:
                from lxml import etree
                d = etree.SubElement(ln, qn("a:prstDash"))
            d.set("val", dash)
        except Exception:
            pass
    connector.shadow.inherit = False
    return connector


def add_divider(slide, x, y, w, color: RGBColor = C_HAIR, width_pt: float = 0.5):
    return add_line(slide, x, y, x + w, y, color, width_pt=width_pt)


# ---------------------------------------------------------------------------
# Layout helper — vertical justify
# ---------------------------------------------------------------------------

def vstack(y_top, y_bottom, block_heights: list, min_gap=GAP_CARD) -> list:
    """Return y positions distributing residual space evenly into gaps.

    Eliminates the 'all content packed at top, dead space at bottom'
    anti-pattern: sum block heights, divide leftover into the (n-1) gaps.
    Falls back to min_gap when blocks overflow the area."""
    n = len(block_heights)
    if n == 0:
        return []
    total = sum(int(h) for h in block_heights)
    avail = int(y_bottom) - int(y_top)
    if n == 1:
        return [int(y_top) + max(0, (avail - total) // 2)]
    residual = avail - total
    gap = max(int(min_gap), residual // (n - 1)) if residual > 0 else int(min_gap)
    # Cap gap to avoid absurd spreads on near-empty slides
    gap = min(gap, int(Inches(0.72)))
    used = total + gap * (n - 1)
    y = int(y_top) + max(0, (avail - used) // 2)
    ys = []
    for h in block_heights:
        ys.append(y)
        y += int(h) + gap
    return ys


# ---------------------------------------------------------------------------
# Gradient (legacy — kept for cover band use only; NOT used by header v2)
# ---------------------------------------------------------------------------

def _add_gradient_rect(slide, x, y, w, h, color_top: RGBColor,
                       color_bottom: RGBColor, angle: int = 90):
    from lxml import etree
    shape = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    shape.line.fill.background()
    sp_pr = shape._element.find(qn("p:spPr"))
    if sp_pr is None:
        sp_pr = etree.SubElement(shape._element, qn("p:spPr"))
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        for child in list(sp_pr.findall(qn(tag))):
            sp_pr.remove(child)
    gf = etree.Element(qn("a:gradFill"))
    gsl = etree.SubElement(gf, qn("a:gsLst"))
    gs1 = etree.SubElement(gsl, qn("a:gs"), attrib={"pos": "0"})
    etree.SubElement(gs1, qn("a:srgbClr"), attrib={"val": str(color_top)})
    gs2 = etree.SubElement(gsl, qn("a:gs"), attrib={"pos": "100000"})
    etree.SubElement(gs2, qn("a:srgbClr"), attrib={"val": str(color_bottom)})
    etree.SubElement(gf, qn("a:lin"), attrib={"ang": str(angle * 60000), "scaled": "1"})
    prst_geom = sp_pr.find(qn("a:prstGeom"))
    if prst_geom is not None:
        idx = list(sp_pr).index(prst_geom) + 1
        sp_pr.insert(idx, gf)
    else:
        sp_pr.insert(0, gf)
    for style_el in list(shape._element.findall(qn("p:style"))):
        shape._element.remove(style_el)
    return shape


# ---------------------------------------------------------------------------
# Chrome: header / footer (v2 — destructive rewrite)
# ---------------------------------------------------------------------------

_GLYPH_PREFIX = re.compile(r"^[◆★■▼●♦]\s*")


def add_header_bar(slide, title: str, logo_path: Optional[Path] = None,
                   eyebrow: Optional[str] = None):
    """v2 header: white zone, navy title, content-width navy rule with an
    orange brand tick. Color logo top-right. No filled bar, no gradient."""
    title = _GLYPH_PREFIX.sub("", str(title))

    if eyebrow:
        add_textbox(slide, MARGIN, Inches(0.22),
                    _CONTENT_W - Inches(2.0), Inches(0.20),
                    eyebrow,
                    font_name=FONT_BODY, font_size_pt=SIZE_CAPTION,
                    font_color=C_ORANGE, bold=True, tracking_pt=1.2)
        title_y = Inches(0.44)
    else:
        title_y = Inches(0.34)

    add_textbox(slide, MARGIN, title_y,
                _CONTENT_W - Inches(2.0), Inches(0.44),
                title,
                font_name=FONT_BLACK, font_size_pt=SIZE_H1,
                font_color=C_NAVY, bold=True)

    # Color logo top-right on white canvas
    if logo_path and Path(logo_path).exists():
        logo_h = Inches(0.40)
        logo_w = Inches(1.30)
        try:
            add_image_contain(slide,
                              SLIDE_W - MARGIN - logo_w, Inches(0.28),
                              logo_w, logo_h, Path(logo_path))
        except Exception:
            pass

    # Rule: 0.75pt navy across content width + orange tick overlay
    rule_y = Inches(1.00)
    add_line(slide, MARGIN, rule_y, SLIDE_W - MARGIN, rule_y, C_NAVY,
             width_pt=0.75)
    add_rect(slide, MARGIN, rule_y - Inches(0.03), Inches(1.2), Inches(0.06),
             C_ORANGE)


def add_footer(slide, text: str = "株式会社オルテナジー  |  https://altenergy.co.jp/",
               page: Optional[int] = None, total: Optional[int] = None):
    """v2 footer: hairline + 8pt faint text + orange tick. No filled band."""
    add_line(slide, MARGIN, FOOTER_RULE_Y, SLIDE_W - MARGIN, FOOTER_RULE_Y,
             C_HAIR, width_pt=0.5)
    add_textbox(slide, MARGIN, FOOTER_RULE_Y + Inches(0.05),
                Inches(5.0), Inches(0.18),
                text,
                font_name=FONT_BODY, font_size_pt=SIZE_SMALL,
                font_color=C_FAINT)
    tick = Inches(0.08)
    add_rect(slide, SLIDE_W - MARGIN - tick, FOOTER_RULE_Y + Inches(0.07),
             tick, tick, C_ORANGE)
    if page is not None:
        label = f"{page:02d} / {total:02d}" if total else f"{page:02d}"
        add_textbox(slide, SLIDE_W - MARGIN - Inches(1.25),
                    FOOTER_RULE_Y + Inches(0.05),
                    Inches(1.05), Inches(0.18),
                    label,
                    font_name=FONT_BODY, font_size_pt=SIZE_SMALL,
                    font_color=C_FAINT, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Section header (v2 — orange square marker, glyph prefixes stripped)
# ---------------------------------------------------------------------------

def add_section_header(slide, x, y, w, text: str, font_size_pt: float = SIZE_H2):
    """v2: 0.10in orange square + bold heading.
    Legacy '◆ ' / '★ ' prefixes are stripped automatically."""
    text = _GLYPH_PREFIX.sub("", str(text))
    sq = Inches(0.10)
    add_rect(slide, x, y + Inches(0.07), sq, sq, C_ORANGE)
    add_textbox(slide,
                x + sq + Inches(0.10), y,
                w - sq - Inches(0.10), Inches(0.30),
                text,
                font_name=FONT_BLACK, font_size_pt=font_size_pt,
                font_color=C_DARK, bold=True)


def add_section_header_v2(slide, x, y, w, text: str,
                          subtitle: Optional[str] = None,
                          font_size_pt: float = SIZE_H2,
                          underline: bool = False):
    """Compatibility wrapper over the v2 section header. Returns next y."""
    add_section_header(slide, x, y, w, text, font_size_pt=font_size_pt)
    next_y = y + Inches(0.32)
    if subtitle:
        add_textbox(slide, x + Inches(0.20), next_y, w - Inches(0.20),
                    Inches(0.22),
                    subtitle, font_name=FONT_BODY, font_size_pt=SIZE_CAPTION,
                    font_color=C_SUB)
        next_y += Inches(0.22)
    if underline:
        add_divider(slide, x, next_y + Inches(0.04), w)
    return next_y


# ---------------------------------------------------------------------------
# KPI card (v2 — white card, hairline border, single 28pt number size)
# ---------------------------------------------------------------------------

def add_kpi_card(slide, x, y, w, h,
                 number: str, unit: str, label: str,
                 bg_color: RGBColor = None,          # ignored (v2: always white)
                 number_size_pt: float = None,       # ignored (v2: fixed 28pt)
                 accent_color: Optional[RGBColor] = None,
                 accent_bar: bool = True):
    """v2 KPI card. White card + 0.75pt hairline + inset orange accent bar.
    Number is ALWAYS 28pt orange with the unit on the same baseline.
    bg_color / number_size_pt are accepted for API compatibility but ignored
    so every KPI in the deck renders identically."""
    if accent_color is None:
        accent_color = C_ORANGE

    add_rounded_rect(slide, x, y, w, h, C_WHITE,
                     radius_pt=CARD_RADIUS_PT,
                     border_color=C_HAIR, border_pt=0.75)
    if accent_bar:
        add_rect(slide, x + ACCENT_BAR_INSET, y,
                 w - ACCENT_BAR_INSET * 2, ACCENT_BAR_H, accent_color)

    # Label top-left, 9pt sub
    add_textbox(slide, x + CARD_PAD, y + Inches(0.12),
                w - CARD_PAD * 2, Inches(0.20),
                str(label),
                font_name=FONT_BODY, font_size_pt=SIZE_CAPTION,
                font_color=C_SUB, bold=True)

    # Number + unit, baseline pair, fills the lower area
    add_number_unit(slide, x + CARD_PAD, y + Inches(0.30),
                    w - CARD_PAD * 2, h - Inches(0.44),
                    str(number), str(unit),
                    number_size_pt=SIZE_KPI,
                    number_color=accent_color,
                    align=PP_ALIGN.LEFT)


def add_metric_hero(slide, x, y, w, h,
                    number: str, unit: str, label: str,
                    accent_color: RGBColor = C_ORANGE,
                    bg_color: RGBColor = None,        # ignored — no background
                    number_size_pt: float = SIZE_HERO):
    """v2 hero metric: eyebrow label → 48pt number + unit run → orange rule.
    Sits directly on the white canvas (no card chrome). Use at most twice
    per deck (price hero, summary cumulative)."""
    add_textbox(slide, x, y, w, Inches(0.22),
                str(label),
                font_name=FONT_BODY, font_size_pt=SIZE_CAPTION,
                font_color=C_ORANGE, bold=True, tracking_pt=1.2)
    add_number_unit(slide, x, y + Inches(0.24),
                    w, h - Inches(0.70),
                    str(number), str(unit),
                    number_size_pt=number_size_pt,
                    unit_size_pt=16,
                    number_color=accent_color,
                    align=PP_ALIGN.LEFT)
    rule_w = min(int(Inches(2.0)), int(w))
    add_line(slide, x, y + h - Inches(0.18),
             int(x) + rule_w, y + h - Inches(0.18),
             C_ORANGE, width_pt=0.75)


def add_pill_label(slide, x, y, w, h, text: str,
                   bg_color: RGBColor = None,
                   font_color: RGBColor = C_ORANGE,
                   font_size_pt: float = SIZE_CAPTION):
    """v2 pill: outline style (no fill), fixed adj=50000 for true pill ends."""
    shape = slide.shapes.add_shape(5, int(x), int(y), int(w), int(h))
    shape.fill.background()
    shape.line.color.rgb = font_color
    shape.line.width = Pt(0.75)
    sp_pr = shape.element.find(qn("p:spPr"))
    if sp_pr is not None:
        prstgeom = sp_pr.find(qn("a:prstGeom"))
        if prstgeom is not None:
            av_lst = prstgeom.find(qn("a:avLst"))
            if av_lst is not None:
                for gd in av_lst.findall(qn("a:gd")):
                    if gd.get("name") == "adj":
                        gd.set("fmla", "val 50000")
    shape.shadow.inherit = False
    text_h = Pt(font_size_pt * 1.5)
    add_textbox(slide, x, int(y) + (int(h) - int(text_h)) // 2, w, text_h,
                str(text),
                font_name=FONT_BODY, font_size_pt=font_size_pt,
                font_color=font_color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_number_marker(slide, cx, cy, number: str, diameter=Inches(0.34)):
    """True circle number marker: white fill, 1pt orange ring,
    12pt bold orange numeral centered. Replaces ①② glyphs and
    rounded-square pseudo-circles."""
    x = int(cx) - int(diameter) // 2
    y = int(cy) - int(diameter) // 2
    add_oval(slide, x, y, diameter, diameter,
             fill_color=C_WHITE, border_color=C_ORANGE, border_pt=1.0)
    tb = slide.shapes.add_textbox(x, y, int(diameter), int(diameter))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(number)
    r.font.name = FONT_BLACK
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = C_ORANGE
    return tb


def add_card_with_accent(slide, x, y, w, h,
                         accent_color: RGBColor = C_ORANGE,
                         bg_color: RGBColor = C_WHITE,
                         border: bool = True,
                         accent_position: str = "top"):
    """White card + hairline + inset accent stripe.
    Returns inner content bounds (cx, cy, cw, ch)."""
    add_rounded_rect(slide, x, y, w, h, bg_color,
                     radius_pt=CARD_RADIUS_PT,
                     border_color=C_HAIR if border else None,
                     border_pt=0.75 if border else 0.0)
    if accent_position == "top":
        add_rect(slide, x + ACCENT_BAR_INSET, y,
                 w - ACCENT_BAR_INSET * 2, ACCENT_BAR_H, accent_color)
        return (x + CARD_PAD, y + ACCENT_BAR_H + Inches(0.08),
                w - CARD_PAD * 2, h - ACCENT_BAR_H - Inches(0.16))
    elif accent_position == "left":
        add_rect(slide, x, y + ACCENT_BAR_INSET,
                 Inches(0.05), h - ACCENT_BAR_INSET * 2, accent_color)
        return (x + Inches(0.05) + CARD_PAD, y + CARD_PAD,
                w - Inches(0.05) - CARD_PAD * 2, h - CARD_PAD * 2)
    return (x + CARD_PAD, y + CARD_PAD, w - CARD_PAD * 2, h - CARD_PAD * 2)


# ---------------------------------------------------------------------------
# Table (v2 — "audited figures" style)
# ---------------------------------------------------------------------------

_NUMERIC_RE = re.compile(r"^[¥\\\-▲△+]?[\d,，.]+\s*[%％円万kWhWm²年℃t\-—/]*$")

TABLE_ROW_H = Inches(0.26)

_TCPR_ORDER = ["lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr", "cell3D",
               "noFill", "solidFill", "gradFill", "blipFill", "pattFill",
               "grpFill", "headers", "extLst"]


def _tcpr_sort(tcPr):
    """Re-order tcPr children to satisfy OOXML schema order."""
    def key(el):
        name = el.tag.split('}')[-1]
        try:
            return _TCPR_ORDER.index(name)
        except ValueError:
            return len(_TCPR_ORDER)
    children = sorted(list(tcPr), key=key)
    for el in list(tcPr):
        tcPr.remove(el)
    for el in children:
        tcPr.append(el)


def _set_cell_border(cell, edge: str, color: Optional[RGBColor],
                     width_pt: float = 0.5, dash: Optional[str] = None):
    """Set one border edge of a table cell ('T','B','L','R').
    color=None → explicit no-line."""
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tag = qn(f"a:ln{edge}")
    for el in tcPr.findall(tag):
        tcPr.remove(el)
    ln = etree.SubElement(tcPr, tag)
    ln.set("w", str(int(width_pt * 12700)))
    ln.set("cap", "flat")
    if color is None:
        etree.SubElement(ln, qn("a:noFill"))
    else:
        fill = etree.SubElement(ln, qn("a:solidFill"))
        clr = etree.SubElement(fill, qn("a:srgbClr"))
        clr.set("val", str(color))
        if dash:
            d = etree.SubElement(ln, qn("a:prstDash"))
            d.set("val", dash)
    _tcpr_sort(tcPr)


def _set_cell_bg(cell, color: RGBColor):
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for el in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(el)
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", str(color))
    _tcpr_sort(tcPr)


def add_table(slide, x, y, w, rows_data: list[list],
              col_widths: list,
              header_bg: RGBColor = None,          # ignored (v2: white header)
              row_bg_even: RGBColor = None,        # ignored
              row_bg_odd: RGBColor = None,         # ignored
              font_size_pt: float = SIZE_TABLE,
              highlight_col: Optional[int] = None,
              total_row: Optional[int] = None):
    """v2 table: white header w/ navy bold text + 1pt navy bottom rule,
    warm zebra body rows, horizontal hairlines only, numerals right-aligned.
    Table height = len(rows)*TABLE_ROW_H — use it to advance y safely.
    highlight_col: column index to tint (#FDEDE6, e.g. PPA案 column).
    total_row: row index rendered as total (bold, tint bg, navy top rule)."""
    from pptx.util import Pt as _Pt
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, int(x), int(y),
                                       int(w), int(TABLE_ROW_H) * n_rows)
    tbl = tbl_shape.table

    # Kill theme banding
    try:
        tblPr = tbl._tbl.find(qn("a:tblPr"))
        if tblPr is not None:
            tblPr.set("firstRow", "0")
            tblPr.set("bandRow", "0")
    except Exception:
        pass

    for c, cw in enumerate(col_widths):
        tbl.columns[c].width = int(cw)
    for r in range(n_rows):
        tbl.rows[r].height = int(TABLE_ROW_H)

    for r, row in enumerate(rows_data):
        is_header = (r == 0)
        is_total = (total_row is not None and r == total_row)
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            text = str(cell_text) if cell_text is not None else ""
            cell.text = text
            cell.margin_left = _Pt(4)
            cell.margin_right = _Pt(4)
            cell.margin_top = _Pt(1)
            cell.margin_bottom = _Pt(1)
            try:
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

            numeric = (bool(_NUMERIC_RE.match(text.replace(" ", "")))
                       and not is_header and text not in ("", "—", "-"))
            for para in cell.text_frame.paragraphs:
                if is_header:
                    para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                else:
                    para.alignment = (PP_ALIGN.RIGHT if numeric and c > 0
                                      else (PP_ALIGN.LEFT if c == 0
                                            else PP_ALIGN.CENTER))
                for run in para.runs:
                    run.font.name = FONT_BODY
                    run.font.size = _Pt(font_size_pt)
                    run.font.bold = is_header or is_total
                    if is_header:
                        run.font.color.rgb = C_NAVY
                    elif is_total and numeric:
                        run.font.color.rgb = C_ORANGE
                    elif numeric and text.startswith(("▲", "△", "-")):
                        run.font.color.rgb = C_ORANGE_DARK
                    else:
                        run.font.color.rgb = C_DARK

            # Fills
            if is_header:
                _set_cell_bg(cell, C_WHITE)
            elif is_total or (highlight_col is not None and c == highlight_col):
                _set_cell_bg(cell, C_TINT)
            elif r % 2 == 0:
                _set_cell_bg(cell, C_ZEBRA)
            else:
                _set_cell_bg(cell, C_WHITE)

            # Borders: horizontal only
            try:
                _set_cell_border(cell, "L", None)
                _set_cell_border(cell, "R", None)
                if is_header:
                    _set_cell_border(cell, "T", None)
                    _set_cell_border(cell, "B", C_NAVY, 1.0)
                else:
                    if is_total:
                        _set_cell_border(cell, "T", C_NAVY, 0.75)
                    else:
                        _set_cell_border(cell, "T", None)
                    _set_cell_border(cell, "B", C_HAIR, 0.5)
            except Exception:
                pass  # borders are cosmetic; never break generation

    return tbl


# ---------------------------------------------------------------------------
# Chart styling (v2 — grayscale-safe double encoding)
# ---------------------------------------------------------------------------

def style_chart_base(chart, font_size_pt: float = SIZE_CAPTION):
    """Frameless chart, dashed warm gridlines, 9pt #666 axes/legend."""
    from pptx.util import Pt as _Pt
    try:
        chart.font.size = _Pt(font_size_pt)
        chart.font.name = FONT_BODY
        chart.font.color.rgb = C_SUB
    except Exception:
        pass
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        gl = va.major_gridlines.format.line
        gl.color.rgb = C_GRID_LINE
        gl.width = _Pt(0.5)
        try:
            ln = gl._get_or_add_ln()
            d = ln.find(qn("a:prstDash"))
            if d is None:
                from lxml import etree
                d = etree.SubElement(ln, qn("a:prstDash"))
            d.set("val", "dash")
        except Exception:
            pass
        va.format.line.fill.background()
    except Exception:
        pass
    try:
        ca = chart.category_axis
        ca.format.line.color.rgb = C_HAIR
        ca.tick_labels.font.size = _Pt(font_size_pt)
    except Exception:
        pass


def style_series_before(series):
    """'Before / current' series: 1.5pt gray DASHED — never orange.
    Grayscale-print double encoding: color AND dash differ."""
    from pptx.util import Pt as _Pt
    series.format.line.color.rgb = C_CHART_GRAY
    series.format.line.width = _Pt(1.5)
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        series.format.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    except Exception:
        pass
    try:
        series.smooth = False
    except Exception:
        pass


def style_series_after(series):
    """'After / proposal' series: 2.25pt orange SOLID."""
    from pptx.util import Pt as _Pt
    series.format.line.color.rgb = C_ORANGE
    series.format.line.width = _Pt(2.25)
    try:
        series.smooth = False
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------

def fmt_yen(value, unit: str = "円") -> str:
    if value is None:
        return "—"
    try:
        v = float(value)
        if abs(v) >= 1_0000_0000:
            return f"{v / 1_0000_0000:.1f}億{unit}"
        if abs(v) >= 10_000:
            return f"{v / 10_000:,.0f}万{unit}"
        return f"{v:,.0f}{unit}"
    except (TypeError, ValueError):
        return str(value)


def fmt_num(value, decimals: int = 1, suffix: str = "") -> str:
    if value is None:
        return "—"
    try:
        return f"{float(value):,.{decimals}f}{suffix}"
    except (TypeError, ValueError):
        return str(value)
