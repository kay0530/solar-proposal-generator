"""
generator.py - PPTX proposal generation engine

Dynamically imports slide modules (PP0→pp0, PP4A→pp4a, NEW_summary→new_summary)
and assembles them into a single PPTX file.

Supports:
  - Profile-based slide lists (from composition_profiles.yaml)
  - Custom slide list overrides (for irregular cases like FIP combos)
  - Mixed PPA/EPC/NEW slides in a single deck
"""

from __future__ import annotations

import importlib
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches

from proposal_generator.utils import SLIDE_W, SLIDE_H

# ---------------------------------------------------------------------------
# Slide module resolution
# ---------------------------------------------------------------------------

# Base directories for slide modules
_BASE_DIR = Path(__file__).parent
_SLIDE_DIRS = {
    "ppa": _BASE_DIR / "slides" / "ppa",
    "epc": _BASE_DIR / "slides" / "epc",
    "new": _BASE_DIR / "slides" / "new",
}

LOGO_PATH = _BASE_DIR / "logo.png"

# Required slides that should always be present (warnings if missing)
REQUIRED_SLIDES = {"PP0", "EP0"}  # At least one cover slide


def _slide_id_to_module_name(slide_id: str) -> str:
    """Convert slide ID to Python module filename (without .py).

    PP0  → pp0
    PP4A → pp4a
    PP8A → pp8a
    EP0  → ep0
    NEW_summary → new_summary
    NEW_fip → new_fip
    """
    return slide_id.lower()


def _resolve_slide_module(slide_id: str):
    """Dynamically import and return the slide module for a given ID.

    Search order:
      1. slides/ppa/ (PP* slides)
      2. slides/epc/ (EP* slides)
      3. slides/new/ (NEW_* slides)

    Returns the module or None if not found.
    """
    module_name = _slide_id_to_module_name(slide_id)

    # Determine which subdirectory to search first based on prefix
    if slide_id.upper().startswith("PP"):
        search_order = ["ppa", "new", "epc"]
    elif slide_id.upper().startswith("EP"):
        search_order = ["epc", "new", "ppa"]
    elif slide_id.upper().startswith("NEW"):
        search_order = ["new", "ppa", "epc"]
    else:
        search_order = ["new", "ppa", "epc"]

    for subdir in search_order:
        module_path = _SLIDE_DIRS[subdir] / f"{module_name}.py"
        if module_path.exists():
            # Build the fully qualified module name
            fq_name = f"proposal_generator.slides.{subdir}.{module_name}"
            try:
                return importlib.import_module(fq_name)
            except ImportError:
                continue

    return None


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

def validate_slide_list(slide_ids: list[str], catalog: dict | None = None) -> list[str]:
    """Validate a slide list and return warnings.

    Args:
        slide_ids: List of slide IDs to validate.
        catalog: Optional slide catalog dict for title lookups.

    Returns:
        List of warning messages (empty if all OK).
    """
    warnings = []

    if not slide_ids:
        warnings.append("スライドが選択されていません")
        return warnings

    # Check for at least one cover slide
    has_cover = any(sid in REQUIRED_SLIDES for sid in slide_ids)
    if not has_cover:
        warnings.append("表紙スライド（PP0 または EP0）が含まれていません")

    # Check for duplicate slides
    seen = set()
    for sid in slide_ids:
        if sid in seen:
            warnings.append(f"重複スライド: {sid}")
        seen.add(sid)

    # Check that all slides have modules
    for sid in slide_ids:
        mod = _resolve_slide_module(sid)
        if mod is None:
            warnings.append(f"スライドモジュールが見つかりません: {sid}")

    return warnings


# ---------------------------------------------------------------------------
# Main generation function
# ---------------------------------------------------------------------------

def generate_proposal(
    slide_ids: list[str],
    data: dict,
    output_path: Path,
    template_path: Optional[Path] = None,
    logo_path: Optional[Path] = None,
) -> dict:
    """Generate a PPTX proposal from a list of slide IDs and customer data.

    This function accepts any combination of slide IDs (PPA, EPC, NEW),
    enabling custom compositions for irregular cases like FIP + self-consumption.

    Args:
        slide_ids:     Ordered list of slide IDs to include.
        data:          Customer data dict (passed to each slide's generate()).
        output_path:   Where to save the generated PPTX.
        template_path: Optional PPTX template file.
        logo_path:     Optional logo image path.

    Returns:
        Dict with generation results:
          - slides_generated: number of slides successfully generated
          - slides_skipped: list of slide IDs that failed
          - warnings: list of warning messages
    """
    if logo_path is None:
        logo_path = LOGO_PATH if LOGO_PATH.exists() else None

    # Create presentation
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()

    # Set slide dimensions to A4 landscape
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    results = {
        "slides_generated": 0,
        "slides_skipped": [],
        "warnings": [],
    }

    # Validate
    validation_warnings = validate_slide_list(slide_ids)
    results["warnings"].extend(validation_warnings)

    # Generate each slide
    for slide_id in slide_ids:
        mod = _resolve_slide_module(slide_id)
        if mod is None:
            results["slides_skipped"].append(slide_id)
            continue

        if not hasattr(mod, "generate"):
            results["slides_skipped"].append(slide_id)
            results["warnings"].append(
                f"{slide_id}: generate() 関数が定義されていません"
            )
            continue

        # Add a blank slide
        try:
            layout = prs.slide_layouts[6]  # Blank layout
        except IndexError:
            layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Call the slide's generate function
        try:
            mod.generate(slide, data, logo_path=logo_path)
            results["slides_generated"] += 1
        except Exception as e:
            results["slides_skipped"].append(slide_id)
            results["warnings"].append(f"{slide_id} 生成エラー: {e}")

    # Save
    prs.save(str(output_path))
    return results
