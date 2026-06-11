"""
pp5.py - 設備レイアウト・積載荷重 (Design v2: Institutional Trust Grid)

Two-column layout:
  - Left (cols 0-6): layout image inside a white hairline drawing frame
    with a 9pt caption row below (frame renders even with no image)
  - Right (cols 7-11): 3 KPI mini-cards (roof load / total weight /
    panel count) above the v2 "audited figures" load table
Gracefully degrades when either image or load data is missing.
The dark compass rose overlay is kept as-is (repositioned to the new
frame geometry only).
"""
from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from proposal_generator.utils import (
    CARD_PAD, CONTENT_BOTTOM, CONTENT_TOP, C_DARK, C_HAIR, C_PANEL, C_SUB,
    C_WHITE, FONT_BLACK, FONT_BODY, GAP_BLOCK, GAP_IN_CARD, MARGIN,
    SIZE_BODY, SIZE_CAPTION, SIZE_LEAD, SLIDE_W, TABLE_ROW_H,
    add_footer, add_header_bar, add_image_contain, add_kpi_card,
    add_rounded_rect, add_section_header, add_table, add_textbox,
    grid_w, grid_x, vstack,
)

TITLE = "設備レイアウト・積載荷重"
EYEBROW = "03｜効果シミュレーション"

SECTION_HEADER_H = Inches(0.45)   # section header + breathing room
CAPTION_H = Inches(0.22)          # caption row under the drawing frame


def generate(slide, data: dict, logo_path: Path = None) -> None:
    """Render PP5 (equipment layout & load calculation)."""
    add_header_bar(slide, TITLE, logo_path, eyebrow=EYEBROW)

    y = CONTENT_TOP
    has_image = bool(data.get("layout_image_path"))
    has_load = bool(data.get("load_calc"))

    if has_image and has_load:
        _render_two_column(slide, data, y)
    elif has_image:
        _render_image_only(slide, data, y)
    elif has_load:
        _render_load_only(slide, data, y)
    else:
        _render_fallback(slide, data, y)

    # Overlay compass indicator when angle is specified
    compass_angle = data.get("compass_angle")
    if compass_angle is not None and has_image:
        # Frame top-right corner in the new v2 geometry
        if has_load:
            _img_right = grid_x(0) + grid_w(7)
        else:
            _img_right = SLIDE_W - MARGIN
        _img_top = CONTENT_TOP + SECTION_HEADER_H
        _render_compass_indicator(slide, compass_angle, _img_right, _img_top)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Layout branches
# ---------------------------------------------------------------------------

def _render_two_column(slide, data: dict, y) -> None:
    """Layout image frame (cols 0-6) + load calc panel (cols 7-11)."""
    left_x = grid_x(0)
    left_w = grid_w(7)
    right_x = grid_x(7)
    right_w = grid_w(5)

    add_section_header(slide, left_x, y, left_w, "設備レイアウト図")
    _render_image_frame(slide, data, left_x, y + SECTION_HEADER_H,
                        left_w, CONTENT_BOTTOM)

    add_section_header(slide, right_x, y, right_w, "積載荷重計算")
    _render_load_panel(slide, data, right_x, y + SECTION_HEADER_H,
                       right_w, CONTENT_BOTTOM)


def _render_image_only(slide, data: dict, y) -> None:
    """Full-width layout image frame; system info shown in the caption."""
    full_w = SLIDE_W - MARGIN * 2
    add_section_header(slide, MARGIN, y, full_w, "設備レイアウト図")
    _render_image_frame(slide, data, MARGIN, y + SECTION_HEADER_H,
                        full_w, CONTENT_BOTTOM)


def _render_load_only(slide, data: dict, y) -> None:
    """Centered load calc panel (no image) with system info band on top."""
    full_w = SLIDE_W - MARGIN * 2
    band_h = Inches(0.50)
    _render_system_info_band(slide, data, MARGIN, y, full_w, band_h)

    y2 = y + band_h + GAP_BLOCK
    add_section_header(slide, MARGIN, y2, full_w, "積載荷重計算結果")

    table_w = min(int(full_w), int(Inches(7.0)))
    table_x = MARGIN + (int(full_w) - table_w) // 2
    _render_load_panel(slide, data, table_x, y2 + SECTION_HEADER_H,
                       table_w, CONTENT_BOTTOM)


def _render_fallback(slide, data: dict, y) -> None:
    """Fallback: system info band + upload prompt, vertically justified."""
    full_w = SLIDE_W - MARGIN * 2
    band_h = Inches(0.50)
    msg_h = Inches(0.40)
    ys = vstack(y, CONTENT_BOTTOM, [band_h, msg_h])
    _render_system_info_band(slide, data, MARGIN, ys[0], full_w, band_h)
    add_textbox(slide, MARGIN, ys[1], full_w, msg_h,
                "レイアウト画像または積載荷重計算表をアップロードしてください。",
                font_name=FONT_BODY, font_size_pt=SIZE_LEAD,
                font_color=C_SUB, align=PP_ALIGN.CENTER,
                line_spacing=1.35)


# ---------------------------------------------------------------------------
# Components
# ---------------------------------------------------------------------------

def _system_info_text(data: dict) -> str:
    """Compact system spec line (PV / panels / PCS / battery)."""
    panel_kw = data.get("panel_total_kw", data.get("system_capacity_kw", 0)) or 0
    panel_count = data.get("panel_total_count", data.get("panel_count", 0)) or 0
    pcs_kw = data.get("pcs_total_kw", data.get("pcs_output_kw", 0)) or 0
    battery_kwh = data.get("battery_total_kwh", data.get("battery_kwh", 0)) or 0

    items = []
    if panel_kw:
        items.append(f"PV出力: {panel_kw:,.2f} kW")
    if panel_count:
        items.append(f"パネル枚数: {panel_count:,}枚")
    if pcs_kw:
        items.append(f"PCS出力: {pcs_kw:,.1f} kW")
    if battery_kwh:
        items.append(f"蓄電池: {battery_kwh:,.1f} kWh")
    return "　｜　".join(items) if items else "設備情報未入力"


def _render_system_info_band(slide, data: dict, x, y, w, h) -> None:
    """v2 system info band: warm panel fill + centered bold body text."""
    add_rounded_rect(slide, x, y, w, h, C_PANEL)
    add_textbox(slide, x + CARD_PAD, y, w - CARD_PAD * 2, h,
                _system_info_text(data),
                font_name=FONT_BLACK, font_size_pt=SIZE_BODY,
                font_color=C_DARK, bold=True, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def _render_image_frame(slide, data: dict, x, y, w, y_bottom) -> None:
    """White drawing frame (0.75pt hairline) + contained image + caption."""
    caption_y = int(y_bottom) - int(CAPTION_H)
    frame_h = caption_y - int(GAP_IN_CARD) - int(y)

    add_rounded_rect(slide, x, y, w, frame_h, C_WHITE,
                     border_color=C_HAIR, border_pt=0.75)

    raw_path = data.get("layout_image_path")
    img_path = Path(raw_path) if raw_path else None
    if img_path and img_path.exists():
        add_image_contain(slide, int(x) + int(CARD_PAD), int(y) + int(CARD_PAD),
                          int(w) - int(CARD_PAD) * 2,
                          frame_h - int(CARD_PAD) * 2, img_path)
    else:
        add_textbox(slide, x, int(y) + frame_h // 2 - int(Inches(0.15)),
                    w, Inches(0.30),
                    "レイアウト画像なし",
                    font_name=FONT_BODY, font_size_pt=SIZE_BODY,
                    font_color=C_SUB, align=PP_ALIGN.CENTER)

    add_textbox(slide, x, caption_y, w, CAPTION_H,
                _system_info_text(data),
                font_name=FONT_BODY, font_size_pt=SIZE_CAPTION,
                font_color=C_SUB)


def _render_load_panel(slide, data: dict, x, y_top, w, y_bottom) -> None:
    """KPI mini-card row + v2 load table, vertically justified."""
    lc = data.get("load_calc") or {}
    if not lc:
        return

    rows = _load_rows(lc)
    kpi_h = Inches(1.00)
    table_h = int(TABLE_ROW_H) * len(rows)

    ys = vstack(y_top, y_bottom, [kpi_h, table_h])
    _render_kpi_row(slide, lc, x, ys[0], w, kpi_h)
    add_table(slide, x, ys[1], w, rows,
              col_widths=[int(w) * 58 // 100, int(w) * 42 // 100],
              font_size_pt=9, total_row=7)


def _render_kpi_row(slide, lc: dict, x, y, w, h) -> None:
    """3 KPI mini-cards: roof load / total weight / panel count."""
    roof_load = float(lc.get("load_per_roof_area", 0) or 0)
    total_weight = float(lc.get("total_weight_kg", 0) or 0)
    panel_count = int(lc.get("panel_count", 0) or 0)

    cards = [
        (f"{roof_load:,.1f}" if roof_load else "—", "kg/㎡", "対屋根面積"),
        (f"{total_weight / 1000:,.1f}" if total_weight else "—", "t", "総重量"),
        (f"{panel_count:,}" if panel_count else "—", "枚", "パネル枚数"),
    ]
    weights = (0.38, 0.32, 0.30)
    usable = int(w) - int(GAP_IN_CARD) * 2
    cx = int(x)
    for (num, unit, label), wt in zip(cards, weights):
        cw = int(usable * wt)
        add_kpi_card(slide, cx, y, cw, h, num, unit, label)
        cx += cw + int(GAP_IN_CARD)


def _load_rows(lc: dict) -> list[list]:
    """Load calc table rows (header + 11 data rows)."""
    return [
        ["項目", "値"],
        ["PV型番", str(lc.get("panel_model") or "—")],
        ["パネル枚数（枚）", f"{int(lc.get('panel_count', 0) or 0):,}"],
        ["パネル単体重量", f"{float(lc.get('panel_unit_weight_kg', 0) or 0):.1f} kg"],
        ["パネル重量（計）", f"{float(lc.get('panel_weight_kg', 0) or 0):,.1f} kg"],
        ["架台重量", f"{float(lc.get('frame_weight_kg', 0) or 0):,.1f} kg"],
        ["配線重量", f"{float(lc.get('wiring_weight_kg', 0) or 0):,.1f} kg"],
        ["総重量", f"{float(lc.get('total_weight_kg', 0) or 0):,.1f} kg"],
        ["パネル面積", f"{float(lc.get('panel_area_m2', 0) or 0):,.1f} m²"],
        ["屋根面積", f"{float(lc.get('roof_area_m2', 0) or 0):,.1f} m²"],
        ["積載荷重（対パネル面積）", f"{float(lc.get('load_per_panel_area', 0) or 0):.2f} kg/m²"],
        ["積載荷重（対屋根面積）", f"{float(lc.get('load_per_roof_area', 0) or 0):.2f} kg/m²"],
    ]


# ---------------------------------------------------------------------------
# Compass direction indicator (rotated by angle) — kept as-is (recent design)
# ---------------------------------------------------------------------------

_ANGLE_TO_LABEL = {
    0: "北", 45: "北東", 90: "東", 135: "南東",
    180: "南", 225: "南西", 270: "西", 315: "北西",
}


def _angle_label(angle: int) -> str:
    """Return a human-readable label for a compass angle."""
    if angle in _ANGLE_TO_LABEL:
        return _ANGLE_TO_LABEL[angle]
    return f"{angle}°"


def _render_compass_indicator(slide, angle: int,
                              img_right=None, img_top=None) -> None:
    """Draw a dark-themed compass rose matching the Streamlit UI style.

    Features:
    - Dark navy filled circle
    - Light grey 8-spoke crosshair (fixed orientation)
    - Cardinal labels: N (orange), E/S/W (light grey) at circle edge
    - Orange arrow (isosceles triangle) rotated by angle
    - Grey counter-arrow pointing opposite direction
    - Bottom label: "XX° - direction"
    """
    from pptx.enum.shapes import MSO_SHAPE
    import math

    box_w = Inches(1.1)
    box_h = Inches(1.3)
    if img_right is not None and img_top is not None:
        box_x = img_right - box_w - Inches(0.1)
        box_y = img_top + Inches(0.1)
    else:
        box_x = SLIDE_W - MARGIN - box_w
        box_y = CONTENT_TOP + Inches(0.05)

    # Compass center + radius
    cx = box_x + box_w / 2
    cy = box_y + Inches(0.55)
    r_outer = Inches(0.45)

    # Colors (dark theme)
    C_COMPASS_BG = RGBColor(0x1E, 0x22, 0x2E)
    C_COMPASS_STROKE = RGBColor(0x3A, 0x3F, 0x4E)
    C_COMPASS_SPOKE = RGBColor(0x55, 0x5A, 0x68)
    C_COMPASS_N = RGBColor(0xE8, 0x49, 0x0F)
    C_COMPASS_CARDINAL = RGBColor(0xB8, 0xBC, 0xC6)
    C_COMPASS_ARROW_BACK = RGBColor(0x7A, 0x7E, 0x8C)

    # 1. Dark filled circle (background)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(cx - r_outer), int(cy - r_outer),
        int(r_outer * 2), int(r_outer * 2),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_COMPASS_BG
    bg.line.color.rgb = C_COMPASS_STROKE
    bg.line.width = Pt(0.75)

    # 2. 8-spoke crosshair (fixed orientation)
    for i in range(8):
        dir_angle = i * 45
        rad = math.radians(dir_angle)
        is_cardinal = (i % 2 == 0)
        r_start = Inches(0.04)
        r_end = r_outer - Inches(0.04)
        line_w = Pt(0.5) if is_cardinal else Pt(0.35)

        x1 = cx + r_start * math.sin(rad)
        y1 = cy - r_start * math.cos(rad)
        x2 = cx + r_end * math.sin(rad)
        y2 = cy - r_end * math.cos(rad)

        connector = slide.shapes.add_connector(
            1, int(x1), int(y1), int(x2), int(y2),
        )
        connector.line.color.rgb = C_COMPASS_SPOKE
        connector.line.width = line_w

    # 3. Cardinal labels (N/E/S/W) at circle edge
    label_r = r_outer - Inches(0.11)
    label_box = Inches(0.16)
    cardinals = [
        ("N", 0, C_COMPASS_N),
        ("E", 90, C_COMPASS_CARDINAL),
        ("S", 180, C_COMPASS_CARDINAL),
        ("W", 270, C_COMPASS_CARDINAL),
    ]
    for lbl, ang, col in cardinals:
        rad = math.radians(ang)
        lx = cx + label_r * math.sin(rad) - label_box / 2
        ly = cy - label_r * math.cos(rad) - label_box / 2
        add_textbox(
            slide, int(lx), int(ly), int(label_box), int(label_box),
            lbl, font_name=FONT_BLACK, font_size_pt=8,
            font_color=col, bold=True, align=PP_ALIGN.CENTER,
        )

    # 4. Arrow (orange) + counter-arrow (grey)
    # Both triangles bbox-centered at (cx, cy) so rotation pivots correctly
    arrow_w = Inches(0.14)
    arrow_h = Inches(0.58)

    # Grey counter-arrow first (so orange draws on top)
    back = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        int(cx - arrow_w / 2), int(cy - arrow_h / 2),
        int(arrow_w), int(arrow_h),
    )
    back.fill.solid()
    back.fill.fore_color.rgb = C_COMPASS_ARROW_BACK
    back.line.fill.background()
    back.rotation = float((angle + 180) % 360)

    # Orange arrow pointing to rotated north
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        int(cx - arrow_w / 2), int(cy - arrow_h / 2),
        int(arrow_w), int(arrow_h),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C_COMPASS_N
    arrow.line.fill.background()
    arrow.rotation = float(angle)

    # 5. Bottom label: "40° - 北東"
    label = _angle_label(angle)
    add_textbox(
        slide, box_x, box_y + Inches(1.05),
        box_w, Inches(0.20),
        f"{angle}° – {label}",
        font_name=FONT_BODY, font_size_pt=9,
        font_color=C_DARK, bold=True, align=PP_ALIGN.CENTER,
    )
