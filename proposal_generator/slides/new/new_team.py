"""
new_team.py - ご支援体制図 (design v2: Institutional Trust Grid)

Placeholder slide showing 5 department cards (営業 / 設計・開発 / 施工管理 /
O&M・保守 / モニタリング) as white cards with top accents and circle
number markers for the flow. Photo placeholders are hairline-framed
circles with 9pt captions; names / roles / duties stay as placeholders.
"""
from __future__ import annotations

from pathlib import Path

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from proposal_generator.utils import (
    CONTENT_BOTTOM, CONTENT_TOP, C_DARK, C_FAINT, C_HAIR, C_NAVY, C_SUB,
    FONT_BLACK, FONT_BODY, GAP_CARD, MARGIN, SIZE_BODY, SIZE_CAPTION,
    SIZE_SMALL, SLIDE_W,
    add_card_with_accent, add_footer, add_header_bar, add_number_marker,
    add_multiline_textbox, add_oval, add_textbox, vstack,
)

TITLE = "ご支援体制図"
EYEBROW = "補足｜ご支援体制"

DEPARTMENTS = [
    ("営業", "Sales"),
    ("設計・開発", "Engineering"),
    ("施工管理", "Construction"),
    ("O&M・保守", "Maintenance"),
    ("モニタリング", "Monitoring"),
]


def _photo_placeholder(slide, x, y, d) -> None:
    """Hairline-framed circle with a 9pt 'PHOTO' caption inside."""
    shape = add_oval(slide, x, y, d, d,
                     fill_color=None, border_color=C_HAIR, border_pt=1.0)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PHOTO"
    run.font.name = FONT_BODY
    run.font.size = Pt(SIZE_CAPTION)
    run.font.color.rgb = C_FAINT


def generate(slide, data: dict, logo_path: Path = None) -> None:
    add_header_bar(slide, TITLE, logo_path, eyebrow=EYEBROW)

    content_w = SLIDE_W - MARGIN * 2

    # ---- Vertical layout: lead + card band + note ----
    lead_h = Inches(0.26)
    card_h = Inches(4.40)
    note_h = Inches(0.22)
    ys = vstack(CONTENT_TOP, CONTENT_BOTTOM, [lead_h, card_h, note_h])

    add_textbox(slide, MARGIN, ys[0], content_w, lead_h,
                "マーケティング～販売～支援～設計～開発～施工～分析～メンテナンス",
                font_size_pt=SIZE_BODY, font_color=C_SUB,
                align=PP_ALIGN.CENTER)

    # ---- 5 department cards ----
    n_cols = len(DEPARTMENTS)
    card_w = (int(content_w) - int(GAP_CARD) * (n_cols - 1)) // n_cols

    for i, (dept_ja, dept_en) in enumerate(DEPARTMENTS):
        x = int(MARGIN) + i * (card_w + int(GAP_CARD))
        cx, cy, cw, ch = add_card_with_accent(slide, x, ys[1], card_w,
                                              card_h, accent_position="top")
        center_x = x + card_w // 2

        # Flow number marker (1..5)
        add_number_marker(slide, center_x, cy + Inches(0.22), str(i + 1),
                          diameter=Inches(0.32))

        # Department name (ja bold + en caption)
        add_multiline_textbox(
            slide, cx, cy + Inches(0.46), cw, Inches(0.52),
            [
                (dept_ja, FONT_BLACK, 12, C_NAVY, True, PP_ALIGN.CENTER),
                (dept_en, FONT_BODY, SIZE_SMALL, C_SUB, False,
                 PP_ALIGN.CENTER),
            ],
            line_spacing=1.2)

        # Photo placeholder: hairline circle + 9pt caption
        photo_d = Inches(1.00)
        _photo_placeholder(slide, center_x - int(photo_d) // 2,
                           cy + Inches(1.10), photo_d)

        # Name / role / duties placeholders
        add_multiline_textbox(
            slide, cx, cy + Inches(2.30), cw, ch - Inches(2.30),
            [
                ("担当者名", FONT_BLACK, 11, C_DARK, True, PP_ALIGN.CENTER),
                ("役職", FONT_BODY, SIZE_CAPTION, C_SUB, False,
                 PP_ALIGN.CENTER),
                ("担当業務の説明をここに記入", FONT_BODY, SIZE_CAPTION,
                 C_SUB, False, PP_ALIGN.CENTER),
            ],
            word_wrap=True, line_spacing=1.5)

    # ---- Note ----
    add_textbox(slide, MARGIN, ys[2], content_w, note_h,
                "※ 詳細は後日記入",
                font_size_pt=SIZE_CAPTION, font_color=C_SUB,
                align=PP_ALIGN.RIGHT)

    add_footer(slide)
